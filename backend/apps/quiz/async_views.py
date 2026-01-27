"""
High-Performance Async Proxy View for Quiz Generation

Implements the Asynchronous Proxy Pattern for quiz endpoints.
"""
import os
import re
import textwrap
from io import BytesIO
from datetime import datetime
from zoneinfo import ZoneInfo

import docx
import PyPDF2
import json
import logging
import httpx
import asyncio
from django.http import JsonResponse, HttpResponse
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_http_methods
from django.utils.decorators import sync_and_async_middleware
from asgiref.sync import sync_to_async
from apps.core.async_client import get_async_client, build_fastapi_headers

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas as rl_canvas
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

logger = logging.getLogger(__name__)


@csrf_exempt
@require_http_methods(["POST"])
async def ajax_extract_text(request):
    """
    Async view to extract text from uploaded files.
    Matches frontend key: 'slide_file'
    Uses asyncio.to_thread to make blocking I/O operations non-blocking.
    """
    if 'slide_file' not in request.FILES:
        return JsonResponse({'error': 'No file uploaded'}, status=400)

    file = request.FILES['slide_file']
    filename = file.name.lower()
    file_ext = os.path.splitext(filename)[1]
    max_size = 10 * 1024 * 1024  # 10MB
    
    if file.size > max_size:
        return JsonResponse({'error': 'File too large (max 10MB)'}, status=400)

    try:
        import asyncio
        
        def extract_text_sync():
            """Synchronous text extraction - runs in thread pool"""
            text = ""
            
            # Read file content into memory for processing
            file.seek(0)
            file_content = file.read()
            
            # PDF extraction
            if file_ext == '.pdf':
                import io
                pdf_file = io.BytesIO(file_content)
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                for page in pdf_reader.pages:
                    text += page.extract_text() or ''
                    
            # DOCX extraction
            elif file_ext == '.docx':
                import io
                doc_file = io.BytesIO(file_content)
                doc = docx.Document(doc_file)
                text = '\n'.join([para.text for para in doc.paragraphs])
                
            # PPTX extraction
            elif file_ext == '.pptx':
                from pptx import Presentation
                import io
                pptx_file = io.BytesIO(file_content)
                prs = Presentation(pptx_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text += shape.text + '\n'
                            
            # TXT extraction
            elif file_ext == '.txt':
                text = file_content.decode('utf-8', errors='ignore')
                
            else:
                raise ValueError(f'Unsupported file format: {file_ext}')
            
            return text.strip()
        
        # Run blocking extraction in thread pool
        text = await asyncio.to_thread(extract_text_sync)
        
        if not text:
            return JsonResponse({'error': 'No text could be extracted from the file.'}, status=400)
        
        # Performance truncation 
        if len(text) > 50000:
            text = text[:49900] + '... [truncated]'
        
        logger.info(f"Text extraction successful! File: {filename} ({round(file.size / (1024 * 1024), 2)} MB, {len(text)} chars)")
        
        # Returns 'text' which frontend sets to extractedText
        return JsonResponse({'text': text})
        
    except ValueError as e:
        logger.warning(f"Unsupported file format: {str(e)}")
        return JsonResponse({'error': str(e)}, status=400)
    except Exception as e:
        logger.error(f"Extraction Error: {str(e)}", exc_info=True)
        return JsonResponse({'error': 'Failed to extract text from file. Please try a different file.'}, status=500)


@csrf_exempt
@require_http_methods(["POST"])
async def generate_quiz_api_async(request):
    """
    High-performance async proxy for quiz generation endpoint.
    
    React-facing endpoint to generate a quiz via FastAPI.
    """
    try:
        # Parse request
        data = json.loads(request.body) if request.body else {}
        
        subject = (data.get("subject") or "General").strip()
        study_text = data.get('extractedText', '').strip()
        num_mcq = data.get("num_mcq") or 7
        num_short = data.get("num_short") or 0
        difficulty = (data.get("difficulty") or "medium").strip().lower()
        
        if not subject:
            return JsonResponse({"error": "Subject is required"}, status=400)
        
        try:
            num_mcq = int(num_mcq)
        except (TypeError, ValueError):
            num_mcq = 7
        
        num_mcq = max(1, min(num_mcq, 20))
        num_short = max(0, min(int(num_short or 0), 10))
        
        # Validate study text
        if not study_text or len(study_text.strip()) < 30:
            return JsonResponse({"error": "Study text must be at least 30 characters"}, status=400)
        
        if len(study_text) > 50000:
            study_text = study_text[:50000]
            logger.warning("Study text truncated to 50,000 characters")
        
        payload = {
            "subject": subject,
            "study_text": study_text.strip(),
            "num_mcq": int(num_mcq),
            "num_short": int(num_short),
            "difficulty": difficulty,
        }
        
        # Forward to FastAPI using async client
        client = get_async_client()
        headers = build_fastapi_headers()
        logger.info(f"Headers sent: {headers}")
        
        fastapi_resp = await client.post(
            "/quiz/",
            json=payload,
            headers=headers,
            timeout=60.0 
        )
        
        if fastapi_resp.status_code != 200:
            logger.warning(f"FastAPI quiz call failed: {fastapi_resp.status_code} {fastapi_resp.text}")
            return JsonResponse(
                {"error": "Quiz service temporarily unavailable"}, 
                status=503
            )
        
        quiz_data = fastapi_resp.json()
        
        # Add metadata for frontend compatibility
        import uuid
        quiz_data['id'] = str(uuid.uuid4())
        # Convert time_limit to integer (minutes) to prevent NaN in frontend timer
        quiz_data['time_limit'] = int(data.get('quiz_time', 10))
        quiz_data['created_at'] = None  # Can be set if storing in DB
        
        return JsonResponse(quiz_data)
        
    except json.JSONDecodeError:
        return JsonResponse({"error": "Invalid JSON"}, status=400)
    except httpx.TimeoutException:
        logger.error("FastAPI quiz request timed out")
        return JsonResponse({"error": "Request timeout"}, status=504)
    except httpx.RequestError as e:
        logger.error(f"FastAPI quiz request error: {e}")
        return JsonResponse({"error": "Service unavailable"}, status=503)
    except Exception as e:
        logger.error(f"Error calling FastAPI quiz endpoint: {e}", exc_info=True)
        return JsonResponse({"error": "Internal server error"}, status=500)


@csrf_exempt
@require_http_methods(["POST"])
async def submit_quiz_api_async(request):
    """
    High-performance async endpoint to submit quiz answers and calculate scores.
    
    For MCQ: Compares answer letters directly
    For Short Answer: Uses LLM to evaluate if answer is correct
    
    Receives user answers and quiz data, calculates score, and returns results.
    """
    try:
        # Parse request
        data = json.loads(request.body) if request.body else {}
        
        quiz_data = data.get("quiz_data")  # Full quiz data with questions
        user_answers = data.get("user_answers", {})  # {question_index: answer}
        quiz_id = data.get("quiz_id")
        
        if not quiz_data:
            return JsonResponse({"error": "Quiz data is required"}, status=400)
        
        # Extract questions from quiz data
        mcq_questions = quiz_data.get("mcq_questions", [])
        short_questions = quiz_data.get("short_questions", [])
        all_questions = mcq_questions + short_questions
        
        if not all_questions:
            return JsonResponse({"error": "No questions found in quiz data"}, status=400)
        
        # Calculate scores
        total_questions = len(all_questions)
        correct_count = 0
        details = []
        
        for idx, question in enumerate(all_questions):
            user_answer = user_answers.get(str(idx), "").strip()
            correct_answer = question.get("answer", "").strip()
            is_correct = False
            reasoning = ""
            
            # For MCQ, compare answer letter (A, B, C, D)
            if question.get("type") == "mcq" or question.get("options"):
                # Normalize answers - handle both "A" and "Option A" formats
                user_letter = user_answer.upper()[0] if user_answer else ""
                correct_letter = correct_answer.upper()[0] if correct_answer else ""
                is_correct = user_letter == correct_letter
                reasoning = "MCQ evaluation: Answer letter matched" if is_correct else "MCQ evaluation: Answer letter did not match"
            else:
                # For short answer, use LLM to evaluate
                question_text = question.get("question", "")
                evaluation = await _evaluate_short_answer(question_text, correct_answer, user_answer)
                is_correct = evaluation.get("is_correct", False)
                reasoning = evaluation.get("reasoning", "Evaluation complete")
            
            if is_correct:
                correct_count += 1
            
            details.append({
                "question_index": idx,
                "question": question.get("question", ""),
                "user_answer": user_answer,
                "correct_answer": correct_answer,
                "is_correct": is_correct,
                "explanation": question.get("explanation", ""),
                "reasoning": reasoning
            })
        
        score_percent = round((correct_count / total_questions) * 100, 1) if total_questions > 0 else 0
        
        results = {
            "quiz_id": quiz_id,
            "subject": quiz_data.get("subject", "Unknown"),
            "difficulty": quiz_data.get("difficulty", "medium"),
            "score": correct_count,
            "total": total_questions,
            "score_percent": score_percent,
            "details": details,
            "submitted_at": None  # Can add timestamp if storing in DB
        }
        
        logger.info(f"Quiz submitted: {correct_count}/{total_questions} correct ({score_percent}%)")
        
        return JsonResponse(results)
        
    except json.JSONDecodeError:
        return JsonResponse({"error": "Invalid JSON"}, status=400)
    except Exception as e:
        logger.error(f"Error processing quiz submission: {e}", exc_info=True)
        return JsonResponse({"error": "Internal server error"}, status=500)


async def _evaluate_short_answer(question_text: str, correct_answer: str, user_answer: str) -> dict:
    """
    Use LLM to evaluate if a short answer is correct.
    
    Returns:
        {
            "is_correct": bool,
            "reasoning": str,
            "score": float (0.0-1.0)
        }
    """
    if not user_answer.strip():
        return {"is_correct": False, "reasoning": "No answer provided", "score": 0.0}
    
    evaluation_prompt = f"""You are an expert quiz evaluator. Evaluate the following student answer to a quiz question.

Question: {question_text}

Expected/Model Answer: {correct_answer}

Student's Answer: {user_answer}

Evaluate if the student's answer is correct. Consider:
1. Factual accuracy
2. Completeness (does it cover key points?)
3. Clarity and relevance

Respond in JSON format ONLY:
{{
  "is_correct": true/false,
  "reasoning": "brief explanation (1-2 sentences)",
  "score": 0.0-1.0
}}

IMPORTANT: Return ONLY the JSON object, nothing else. No markdown formatting, no code blocks."""

    try:
        client = get_async_client()
        headers = build_fastapi_headers()
        
        response = await client.post(
            "/chatbot/",
            json={"message": evaluation_prompt},
            headers=headers,
            timeout=30.0
        )
        
        if response.status_code != 200:
            logger.warning(f"LLM evaluation failed: {response.status_code}")
            # Fallback to string matching if LLM fails
            is_correct = user_answer.lower().strip() == correct_answer.lower().strip()
            return {
                "is_correct": is_correct,
                "reasoning": "String match" if is_correct else "Does not match expected answer",
                "score": 1.0 if is_correct else 0.0
            }
        
        response_data = response.json()
        response_text = response_data.get("response", "")
        
        # Extract JSON from response
        try:
            evaluation = json.loads(response_text)
            return {
                "is_correct": evaluation.get("is_correct", False),
                "reasoning": evaluation.get("reasoning", "Evaluation complete"),
                "score": evaluation.get("score", 0.0)
            }
        except json.JSONDecodeError:
            logger.warning(f"Failed to parse LLM evaluation response: {response_text[:200]}")
            # Fallback to basic comparison
            is_correct = user_answer.lower().strip() == correct_answer.lower().strip()
            return {
                "is_correct": is_correct,
                "reasoning": "Automatic evaluation",
                "score": 1.0 if is_correct else 0.0
            }
    except Exception as e:
        logger.error(f"Error evaluating short answer: {e}")
        # Fallback to basic string matching
        is_correct = user_answer.lower().strip() == correct_answer.lower().strip()
        return {
            "is_correct": is_correct,
            "reasoning": "Fallback evaluation",
            "score": 1.0 if is_correct else 0.0
        }


# --- Download Helpers ---

def _safe_filename(name: str, max_len: int = 180) -> str:
    """Make a string safe for use as a filename."""
    if not name:
        return "Quiz_Results"
    s = str(name).strip()
    s = re.sub(r'\s+', '_', s)
    s = re.sub(r'[\\/:"*?<>|]+', '_', s)
    return s[:max_len] or "Quiz_Results"


def _format_quiz_text(results: dict) -> str:
    """Format quiz results as plain text."""
    score = results.get('score', 0)
    total = results.get('total', 0)
    score_percent = results.get('score_percent', 0)
    details = results.get('details', [])
    subject = results.get('subject', 'Quiz')
    difficulty = results.get('difficulty', 'unknown')
    
    tz = ZoneInfo('Africa/Accra')
    timestamp = datetime.now(tz).strftime('%Y-%m-%d %H:%M:%S %Z')
    
    lines = [
        'Lamla AI - Quiz Results',
        '=' * 70,
        f"Subject: {subject}",
        f"Difficulty: {difficulty}",
        f"Generated: {timestamp}",
        f"Score: {score}/{total} ({score_percent:.1f}%)",
        '=' * 70,
        ''
    ]
    
    lines.append('DETAILED ANSWER REVIEW')
    lines.append('-' * 70)
    lines.append('')
    
    for idx, detail in enumerate(details, 1):
        lines.append(f"Q{idx}. {detail.get('question', '')}")
        lines.append(f"Your Answer: {detail.get('user_answer') or '(Unanswered)'}")
        lines.append(f"Correct Answer: {detail.get('correct_answer', '')}")
        status = '✓ CORRECT' if detail.get('is_correct') else '✗ INCORRECT'
        lines.append(f"Status: {status}")
        
        if detail.get('reasoning'):
            lines.append(f"Evaluation: {detail.get('reasoning')}")
        if detail.get('explanation'):
            lines.append(f"Explanation: {detail.get('explanation')}")
        lines.append('')
    
    return '\n'.join(lines)


@csrf_exempt
@require_http_methods(["POST"])
def download_quiz_results(request):
    """Download quiz results as PDF, DOCX, or TXT."""
    try:
        data = json.loads(request.body) if request.body else {}
        results = data.get('results', {})
        file_format = data.get('format', 'pdf').lower()
        
        if not results:
            return JsonResponse({"error": "No results data provided"}, status=400)
        
        subject = results.get('subject', 'Quiz')
        safe_filename = _safe_filename(subject)
        tz = ZoneInfo('Africa/Accra')
        timestamp = datetime.now(tz).strftime('%Y%m%d_%H%M%S')
        filename_base = f"{safe_filename}_Quiz_{timestamp}"
        
        text_content = _format_quiz_text(results)
        
        if file_format == 'pdf':
            if not HAS_REPORTLAB:
                return JsonResponse(
                    {"error": "PDF generation not available. Please use TXT or DOCX format."},
                    status=400
                )
            
            buffer = BytesIO()
            p = rl_canvas.Canvas(buffer, pagesize=letter)
            width, height = letter
            y = height - 40
            
            for line in text_content.split('\n'):
                for wrapped_line in textwrap.wrap(line, width=95):
                    p.drawString(40, y, wrapped_line)
                    y -= 16
                    if y < 40:
                        p.showPage()
                        y = height - 40
            
            p.save()
            buffer.seek(0)
            response = HttpResponse(buffer.getvalue(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{filename_base}.pdf"'
            return response
        
        elif file_format == 'docx':
            buffer = BytesIO()
            doc = docx.Document()
            
            for line in text_content.split('\n'):
                doc.add_paragraph(line)
            
            doc.save(buffer)
            buffer.seek(0)
            response = HttpResponse(
                buffer.getvalue(),
                content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            )
            response['Content-Disposition'] = f'attachment; filename="{filename_base}.docx"'
            return response
        
        else:  # TXT
            response = HttpResponse(text_content, content_type='text/plain')
            response['Content-Disposition'] = f'attachment; filename="{filename_base}.txt"'
            return response
    
    except Exception as e:
        logger.error(f"Error generating download: {e}", exc_info=True)
        return JsonResponse({"error": "Failed to generate file"}, status=500)

