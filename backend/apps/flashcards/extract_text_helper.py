import os
import docx
import PyPDF2
import logging
import re
import unicodedata
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_http_methods
 
logger = logging.getLogger(__name__)

ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}
ALLOWED_MIME_TYPES = {
    ".pdf": {"application/pdf"},
    ".docx": {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    },
    ".pptx": {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    },
    ".txt": {"text/plain"},
}


def _clean_extracted_text(raw_text: str) -> str:
    """
    Normalize extraction artifacts from PDF/PPTX/DOCX:
    - remove control chars (e.g. \x01)
    - normalize unicode
    - preserve meaningful newlines for downstream segmentation
    """
    if not raw_text:
        return ""

    text = unicodedata.normalize("NFKC", raw_text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")

    # Remove non-printable control characters except \n and \t.
    text = "".join(
        ch if (ch in "\n\t" or ch.isprintable()) else " "
        for ch in text
    )

    # Normalize tabs and spacing around lines.
    text = text.replace("\t", " ")
    text = "\n".join(re.sub(r"[ ]{2,}", " ", ln).strip() for ln in text.split("\n"))

    # Collapse excessive blank lines.
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return text


@csrf_exempt
@require_http_methods(["POST"])
async def ajax_extract_text(request):
    """
    Async view to extract text from uploaded files.
    Matches frontend key: 'slide_file'
    Uses asyncio.to_thread to make blocking I/O operations non-blocking.
    """
    if 'slide_file' not in request.FILES:
        return JsonResponse({'error': 'No file uploaded'}, status=400)

    file = request.FILES['slide_file']
    filename = os.path.basename((file.name or "").strip().lower())
    file_ext = os.path.splitext(filename)[1]
    content_type = (getattr(file, "content_type", "") or "").lower().strip()
    max_size = 10 * 1024 * 1024  # 10MB

    if file_ext not in ALLOWED_EXTENSIONS:
        return JsonResponse({'error': f'Unsupported file format: {file_ext or "unknown"}'}, status=400)

    allowed_types = ALLOWED_MIME_TYPES.get(file_ext, set())
    if content_type and content_type not in allowed_types and content_type != "application/octet-stream":
        return JsonResponse({'error': 'Unsupported file MIME type for this extension.'}, status=400)
    
    if file.size > max_size:
        return JsonResponse({'error': 'File too large (max 10MB)'}, status=400)

    try:
        import asyncio
        
        def extract_text_sync():
            """Synchronous text extraction - runs in thread pool"""
            text = ""
            
            # Read file content into memory for processing
            file.seek(0)
            file_content = file.read()
            
            # PDF extraction
            if file_ext == '.pdf':
                import io
                pdf_file = io.BytesIO(file_content)
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                for page in pdf_reader.pages:
                    text += page.extract_text() or ''
                    
            # DOCX extraction
            elif file_ext == '.docx':
                import io
                doc_file = io.BytesIO(file_content)
                doc = docx.Document(doc_file)
                text = '\n'.join([para.text for para in doc.paragraphs])
                
            # PPTX extraction
            elif file_ext == '.pptx':
                from pptx import Presentation
                import io
                pptx_file = io.BytesIO(file_content)
                prs = Presentation(pptx_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text += shape.text + '\n'
                            
            # TXT extraction
            elif file_ext == '.txt':
                text = file_content.decode('utf-8', errors='ignore')
                
            else:
                raise ValueError(f'Unsupported file format: {file_ext}')
            
            return _clean_extracted_text(text)
        
        # Run blocking extraction in thread pool
        text = await asyncio.to_thread(extract_text_sync)
        
        if not text:
            return JsonResponse({'error': 'No text could be extracted from the file.'}, status=400)
        
        # Performance truncation 
        if len(text) > 30000:
            text = text[:29900] + '... [truncated]'
        
        logger.info(f"Text extraction successful! File: {filename} ({round(file.size / (1024 * 1024), 2)} MB, {len(text)} chars)")
        
        # Returns 'text' which frontend sets to extractedText
        return JsonResponse({'text': text})
        
    except ValueError as e:
        logger.warning(f"Unsupported file format: {str(e)}")
        return JsonResponse({'error': str(e)}, status=400)
    except Exception as e:
        logger.error(f"Extraction Error: {str(e)}", exc_info=True)
        return JsonResponse({'error': 'Failed to extract text from file. Please try a different file.'}, status=500)
